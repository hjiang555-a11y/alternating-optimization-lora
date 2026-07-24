"""Generate detailed weekly group meeting PPT — 22 slides with module-level algo
comparison, FLOPs design breakdown, per-variant explanations.

Week: Jul 23-24, 2026. 6 commits. Design: 16:9 white, dark-blue title bar, CN text.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

W, H = Inches(13.333), Inches(7.5)
DARK  = RGBColor(0x1F, 0x49, 0x7D)
RED   = RGBColor(0xC0, 0x50, 0x4D)
ORANGE= RGBColor(0xED, 0x7D, 0x31)
GREEN = RGBColor(0x9B, 0xBB, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF2, 0xF2, 0xF2)
GRAY  = RGBColor(0x60, 0x60, 0x60)
BLACK = RGBColor(0x20, 0x20, 0x20)
BLUE2 = RGBColor(0xE8, 0xF0, 0xF8)
LBLUE = RGBColor(0xCC, 0xDD, 0xEE)

FIG  = os.path.join(os.path.dirname(__file__),"..","docs","figures")
FIN  = os.path.join(FIG,"final_report")
BIL  = os.path.join(FIG,"bilingual_report")

prs = Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]

# ── helpers ─────────────────────────────────────────────────────────
def tbar(slide, text, sub=None):
    b=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Inches(1.05))
    b.fill.solid();b.fill.fore_color.rgb=DARK;b.line.fill.background()
    tf=b.text_frame;tf.word_wrap=True;tf.margin_left=Inches(0.6);tf.margin_top=Inches(0.15)
    p=tf.paragraphs[0];p.text=text;p.font.size=Pt(26);p.font.color.rgb=WHITE;p.font.bold=True
    if sub:
        p2=tf.add_paragraph();p2.text=sub;p2.font.size=Pt(13);p2.font.color.rgb=LBLUE;p2.font.italic=True
def txt(s,l,t,w,h,text,sz=14,c=BLACK,b=False,a=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    p.text=text;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.alignment=a;return tf
def blt(s,l,t,w,h,items,sz=14,c=BLACK):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=bx.text_frame;tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=it;p.font.size=Pt(sz);p.font.color.rgb=c;p.space_after=Pt(4)
    return tf
def img(s,path,l,t,wd=None,ht=None):
    if not os.path.exists(path):
        txt(s,l,t,3,0.5,f"[Missing:{os.path.basename(path)}]",10,RED);return
    kw={};kw["width"]=Inches(wd) if wd else kw;kw["height"]=Inches(ht) if ht else kw
    if wd and ht: kw={"width":Inches(wd),"height":Inches(ht)}
    elif wd: kw={"width":Inches(wd)}
    elif ht: kw={"height":Inches(ht)}
    s.shapes.add_picture(path,Inches(l),Inches(t),**kw)
def sec(title,sub=""):
    s=prs.slides.add_slide(BLANK)
    bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H)
    bg.fill.solid();bg.fill.fore_color.rgb=DARK;bg.line.fill.background()
    txt(s,1,3.2,11,1.5,title,36,WHITE,True,PP_ALIGN.CENTER)
    if sub: txt(s,1,4.5,11,1,sub,16,LBLUE,a=PP_ALIGN.CENTER)
    return s
def table(s,rows,y0,col_x,col_w,highlight_row=0):
    for r,row in enumerate(rows):
        for c,(text,x,w) in enumerate(zip(row,col_x,col_w)):
            is_hdr=(r==0);is_hl=(r==highlight_row)
            bx=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y0+r*0.42),Inches(w),Inches(0.4))
            bx.fill.solid()
            bx.fill.fore_color.rgb=DARK if is_hdr else (BLUE2 if is_hl else (LGRAY if r%2==0 else WHITE))
            bx.line.fill.background()
            tf=bx.text_frame;tf.word_wrap=True;tf.margin_left=Inches(0.08);tf.margin_top=Inches(0.02)
            p=tf.paragraphs[0];p.text=text;p.font.size=Pt(11 if is_hdr else 10)
            p.font.color.rgb=WHITE if is_hdr else BLACK;p.font.bold=is_hdr or is_hl
def txbox(s,left,top,w,h,lines,sz=11,c=BLACK):
    bx=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(left),Inches(top),Inches(w),Inches(h))
    bx.fill.solid();bx.fill.fore_color.rgb=WHITE;bx.line.color.rgb=RGBColor(0xD0,0xD0,0xD0);bx.line.width=Pt(1)
    tf=bx.text_frame;tf.word_wrap=True;tf.margin_left=Inches(0.08);tf.margin_top=Inches(0.05)
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=ln;p.font.size=Pt(sz);p.font.color.rgb=c;p.space_after=Pt(2)

# ═════════════════════════════════════════════════════════════  S1: Title
s=prs.slides.add_slide(BLANK)
bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H)
bg.fill.solid();bg.fill.fore_color.rgb=DARK;bg.line.fill.background()
txt(s,1,1.3,11,2,"A-SYNC 协议族演化\n与 FLOPs 归一化对比实验",40,WHITE,True,PP_ALIGN.CENTER)
txt(s,1,4.0,11,1,"Alternating Optimization for LLM Post-Training -- 周进展汇报",18,LBLUE,a=PP_ALIGN.CENTER)
txt(s,1,5.2,11,0.6,"2026/07/23-24 | 6 commits | 2 GPU experiments | 3 docs | 7 figures",14,RGBColor(0x88,0x99,0xAA),a=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════  S2: Overview
s=prs.slides.add_slide(BLANK);tbar(s,"本周工作概览","6 commits, 3 major deliverables")
blt(s,0.6,1.4,12,5.5,[
    "1. 双语文档交付: EN+ZH A-SYNC 变体报告 (456KB+529KB PDF)",
    "2. 基础理论文档: 残差放大 (rho=1.08) 因果推导 + A-SYNC 12变体分类与修复策略",
    "3. FLOPs归一化实验: A-SYNC CONSTANT vs AdamW vs LoRA on OPT-125m",
    "",
    "关键数字:",
    "  . 6 commits pushed to github.com/hjiang555-a11y + gingersea",
    "  . 10 张新图表 (3 extra for this PPT)",
    "  . ~7000 字技术文档 (中英双语)",
    "  . A-SYNC CONSTANT 48c: OPT-125m PPL 2246->60.7 | Qwen7B PPL 58.8->7.6",
    "  . LoRA AdamW r=8: PPL 37.3 at only 0.013 TFLOPs (70x less compute!)",
],15)

# ═════════════════════════════════════════════════════════════  S3: Protocol A module breakdown
s=prs.slides.add_slide(BLANK);tbar(s,"Protocol A 模块级算法拆解","每个模块用了什么算法? 为什么会发散?")
blt(s,0.3,1.2,6.3,6,[
    "Protocol A (ASP) = ALS + SGD + Perturb 三阶段交替",
    "",
    "模块1: lm_head (输出投影层, 125M params)",
    "  算法: ALS (Alternating Least Squares)",
    "  操作: 对 lm_head 权重做块内闭式L2求解",
    "  W_new = argmin ||X W^T - Y_target||^2",
    "  特点: 一步到块内最优, 但不感知跨层耦合",
    "  问题: 修改lm_head后所有后续层激活值改变",
    "",
    "模块2: Transformer body (24层attention+FFN)",
    "  算法: SGD (随机梯度下降, lr=1e-4)",
    "  操作: 标准前传+反传, 更新所有body参数",
    "  特点: 试图消化 ALS 引入的分布偏移",
    "  问题: 50步 SGD 远不足以恢复 (需~150+)",
    "",
    "模块3: 全参数空间 (lm_head + body)",
    "  算法: Perturbation (高斯噪声, scale=1e-3)",
    "  操作: 向所有参数注入随机扰动",
    "  目的: 鼓励平坦极小值 (隐式SAM)",
    "  问题: 在深层模型上进一步恶化不稳定",
],11)
# Right side: architecture diagram
img(s,os.path.join(FIN,"fig2_arch.png"),7,1.2,5.8,2.8)
img(s,os.path.join(FIN,"fig1_residual.png"),7,4.2,5.8,3.0)

# ═════════════════════════════════════════════════════════════  S4: Why Protocol A fails
s=prs.slides.add_slide(BLANK);tbar(s,"为什么 Protocol A 在>=28层上发散?","残差放大: 从因果干预到指数级崩塌")
blt(s,0.3,1.2,6,6,[
    "Step 1: ALS 对 lm_head 做硬干预",
    "  do(theta_lm := theta_ALS)  -- Pearl do-calculus",
    "  切断入边, 用L2闭式解直接替换参数",
    "",
    "Step 2: 干预偏差通过残差流传播",
    "  定义 delta = h_ALS - h_original",
    "  每一层: h_{l+1} = h_l + f_l(h_l)",
    "           delta_{l+1} = (I + J_l) . delta_l",
    "  J_l = 层l的雅可比矩阵",
    "",
    "Step 3: 级联放大",
    "  ||delta_L|| = ||delta_0|| . prod(I+J_l)",
    "  = ||delta_0|| . rho^{L-1}",
    "  rho = 1.08 (跨层几何均值, 来自4个模型族拟合)",
    "",
    "Step 4: 临界层数",
    "  1.08^27 = 8.0x  放大",
    "  SGD每周期恢复 ~0.005 PPL",
    "  不对称比: 1600:1 -> 灾难性发散",
    "  L_max = ln(C_recovery/||delta||)/ln(1.08) = 26",
],11)

blt(s,6.8,1.2,6,6,[
    "实验确认 (4个模型族, 8个架构):",
    "",
    "| 模型          | L  | rho^{L-1} | 预测  | 实际 PPL |",
    "| OPT-125m      | 12 | 2.3x      | conv  | 106.9    |",
    "| TinyLlama-1B  | 22 | 5.0x      | conv  | 15.5     |",
    "| Qwen0.5B      | 24 | 5.9x      | 边际  | 18.0     |",
    "| Qwen7B        | 28 | 8.0x      | DIV   | 11/11 DIV|",
    "",
    "关键公式: L_max = ln(C_recovery/||delta||)/ln(rho)",
    "  rho: 每层谱范数 (来自 I+J_l 的李雅普诺夫指数)",
    "  C_recovery: SGD 恢复预算 (eta . mu_min . T_SGD)",
    "  ||delta||: ALS 干预幅度",
    "",
    "为什么 LoRA/Protocol C 安全?",
    "  低秩约束: ||delta_eff|| = (r/d) . ||delta_full||",
    "  r=8, d=768 -> delta 缩小 96x",
    "  总放大: 8.0 * 0.01 = 0.08 << SGD 恢复能力",
],10)

# ═════════════════════════════════════════════════════════════  S5: A-SYNC Innovation
s=prs.slides.add_slide(BLANK);tbar(s,"核心创新: A-SYNC 梯度注入 -- 模块级变化","Protocol A (修改权重) -> A-SYNC (塑形梯度)")
blt(s,0.3,1.2,6.2,5.8,[
    "Protocol A (旧): ALS 权重修改",
    "  lm_head:  ALS闭式解 -> 直接替换 W",
    "  body:     SGD 消化分布偏移",
    "  全参数:   Perturb 噪声注入",
    "  问题:     ALS修改进入前传 -> 残差链放大",
    "",
    "A-SYNC (新): 梯度塑形",
    "  Step 1: ALS 计算 delta = W_als - W_before",
    "  Step 2: 权重立即恢复! lm_head 保持预训练值",
    "  Step 3: SGD 时注入 grad_bias = sync * delta",
    "  body:     SGD 步时 grad += sync_strength * delta",
    "  全参数:   去掉 Perturb 阶段 (7B上扰动的反作用)",
    "",
    "关键差异表:",
    "  |     | Protocol A      | A-SYNC CONSTANT         |",
    "  |lm_head| ALS直接改权重  | ALS算delta不修改权重   |",
    "  |body   | SGD被动消化    | SGD + grad梯度偏置注入  |",
    "  |perturb| 高斯噪声全参数 | 完全移除                |",
    "  |sync   | 无             | 0.05 (恒, 不衰减)       |",
    "  |前传   | ALS权重进入    | 从未看到ALS权重          |",
],11)

blt(s,6.8,1.2,6.2,5.8,[
    "为什么梯度注入绕过了残差放大?",
    "",
    "Protocol A: ALS delta -> weight change",
    "  -> 前传: h_L_perturbed = f(h, W+delta_W)",
    "  -> 残差链: (I+J_L)...(I+J_1).delta -> 指数放大",
    "",
    "A-SYNC: ALS delta -> gradient bias",
    "  -> 前传: h_L = f(h, W_original)  (不变!)",
    "  -> 残差链: 未被触发 (权重未改)",
    "  -> 反传: grad += sync*delta  (方向引导)",
    "  -> body 在弱信号下逐渐自适应",
    "",
    "类比: Protocol A = 直接手术 (高风险)",
    "       A-SYNC  = 物理治疗 (渐进式引导)",
    "",
    "结果: 7B (28L) 首次收敛!",
    "  旧 Protocol A: 11/11 发散",
    "  A-SYNC vanilla 8c: PPL 60.9 -> 25.8",
    "  A-SYNC CONSTANT 48c: PPL 58.8 -> 7.6 (C44收敛)",
],11)

# ═════════════════════════════════════════════════════════════  S6: Variant 1-4
s=prs.slides.add_slide(BLANK);tbar(s,"A-SYNC 变体详解 (1/3): 基础变体","Vanilla / No-Perturb / Cosine / CONSTANT")
y=1.2
txbox(s,0.2,y,6.4,2.8,[
    "变体1: A-SYNC Vanilla (8c) -- 原始概念验证",
    "  脚本: _a_sync_qwen7b.py, _a_sync_8cycle_7b.py",
    "  机制: ALS delta -> 梯度偏置注入, sync=0.05",
    "       指数衰减 (sync *= 0.8/cycle), 含扰动阶段",
    "  7B轨迹: 60.9->48.6->38.7->33.8->30.2->27.9->26.5->25.8",
    "  PPL 25.8 | 缓慢单调衰减 (衰减后期信号消失)",
    "  问题: 扰动噪声稀释ALS信号+指数衰减过早停用",
],9)
txbox(s,6.8,y,6.4,2.8,[
    "变体2: A-SYNC No-Perturb (8c) -- 去掉扰动",
    "  脚本: _a_sync_noperturb_7b.py",
    "  机制: 与Vanilla完全相同, 仅移除扰动阶段",
    "  动机: 0.5B上扰动导致 7.8->23.9 PPL 反弹",
    "  7B轨迹: 60.9->38.8->25.0->21.9->19.6->18.3->17.5->16.6",
    "  PPL 16.6 | +9.2 PPL vs Vanilla (扰动纯属有害)",
    "  结论: 7B上扰动阶段不需要, ALS信号已足够",
],9)
y+=3.0
txbox(s,0.2,y,6.4,2.8,[
    "变体3: A-SYNC Cosine (32c) -- 余弦衰减",
    "  脚本: _a_sync_32cycle_7b.py",
    "  机制: sync = 0.05*0.5*(1+cos(pi*t/T)), 相同lr",
    "       扩展到32周期测试收敛渐近线",
    "  7B轨迹: 59.9->33.4(C4)->18.9(C8)->14.7(C12)->",
    "         13.6(C16)->13.3(C20)->13.2(C24-32) 平台!",
    "  PPL 13.2 | C20平台 (衰减杀死了尾部信号)",
    "  关键诊断: 平台启发了 CONSTANT 实验",
],9)
txbox(s,6.8,y,6.4,2.8,[
    "变体4: A-SYNC CONSTANT (24c) -- 无衰减",
    "  脚本: _a_sync_constant_7b.py",
    "  机制: sync=0.05 恒, lr=2e-4 恒, 无任何衰减!",
    "       直接回应 Cosine 平台诊断",
    "  7B轨迹: 61.8->36.0(C4)->18.5(C8)->12.4(C12)->",
    "         10.4(C16)->9.5(C20)->9.0(C24)",
    "  PPL 9.0 | 第一个突破 PPL 10 的变体!",
    "  结论: 衰减策略全部有害 -- 外部衰减和自然收敛竞争",
],9)

# ═════════════════════════════════════════════════════════════  S7: Variant 5-8
s=prs.slides.add_slide(BLANK);tbar(s,"A-SYNC 变体详解 (2/3): 改进变体","CONSTANT 48c / A-CYCLE / EMA / Aligned")
y=1.2
txbox(s,0.2,y,6.4,2.5,[
    "变体5: A-SYNC CONSTANT (48c) -- BEST",
    "  脚本: _a_sync_48cycle_7b.py",
    "  机制: 同24c, 延伸至48周期, sync+lr全恒",
    "  7B轨迹: 58.8->33.0->16.9->12.2->10.6->9.7->",
    "         9.0->8.7->8.4->8.1->8.0->7.8->7.6",
    "  C44收敛: PPL 7.6 | 首个A-SYNC自然收敛平台!",
    "  每cycle Delta: C40+ = -0.2~-0.3 PPL",
    "  最优Protocol A结果, 7.7x vs baseline (58.8)",
],10)
txbox(s,6.8,y,6.4,2.5,[
    "变体6: A-CYCLE Warm Restart (3x8)",
    "  脚本: _a_cycle_7b.py",
    "  机制: 3块x8周期, 每块内Cosine, 块边界重置",
    "  动机: Cosine衰减太快->重启重新注入ALS信号",
    "  7B轨迹: B1C1-8 61.6->42.6->40.4 |",
    "         B2C1-8 30.5->26.1->25.1 | B3C1-8 20.0->16.9->16.5",
    "  PPL 16.5 | 块内衰减浪费4周期/块",
    "  结论: 块内Cosine衰减引入的浪费 > 重启收益",
],10)
y+=2.7
txbox(s,0.2,y,6.4,2.5,[
    "变体7: A-SYNC+EMA -- delta 平滑",
    "  脚本: _a_sync_plus_variants.py, _a_sync_plus_7b.py",
    "  机制: ema_delta = beta*raw + (1-beta)*ema (beta=0.3)",
    "       注入平滑后的delta代替原始每周期delta",
    "  动机: 抑制minibatch高频噪声, 保留一致方向",
    "  0.5B: 所有变体达PPL 5.5 floor (容量底)",
    "  7B: 脚本存在但未运行 (runs/a_sync_plus_7b.json 缺)",
    "  评估: 应自然补充到CONSTANT中, 预期稳定晚期收敛",
],10)
txbox(s,6.8,y,6.4,2.5,[
    "变体8: A-SYNC+Aligned -- 对齐注入",
    "  脚本: _a_sync_plus_variants.py",
    "  机制: proj=(dot(grad,delta)/dot(grad,grad)).clamp(min=0)*grad",
    "       只注入Delta与SGD梯度正对齐的分量",
    "  动机: Delta反向分量产生破坏性干涉",
    "  诊断: 跟踪每周期cos(aligned_delta, grad)",
    "  0.5B: 起始更高(13.8 vs baseline 9.4)但最终同达5.5",
    "  7B: 未运行",
    "  评估: 对齐约束初期降低有效注入量, 晚期可能稳定改善",
],10)

# ═════════════════════════════════════════════════════════════  S8: Variant 9-12
s=prs.slides.add_slide(BLANK);tbar(s,"A-SYNC 变体详解 (3/3): 替代方法","SWA / A-PROBE / A-KD / LARS")
y=1.2
txbox(s,0.2,y,6.4,2.5,[
    "变体9: A-SYNC+SWA Cosine (16c)",
    "  脚本: _a_sync_swa_cosine_7b.py",
    "  机制: SWA=C10起, EMA平均参数; Cosine衰减sync",
    "       最终加载SWA权重做评估",
    "  7B: 标准路径PPL 10.5, SWA模型PPL 13.8",
    "  SWA反而变差 (+3.3 PPL)!",
    "  原因: SWA假设稳态SGD轨迹, A-SYNC是非稳态的",
    "       (ALS delta方向每周变化), 平均不同阶段是坏的",
    "  结论: SWA与A-SYNC根本不相容",
],10)
txbox(s,6.8,y,6.4,2.5,[
    "变体10: A-PROBE (r=64, 16c)",
    "  脚本: _probe_7b.py",
    "  机制: 不修改lm_head, 插入低秩探针(3584->64->3584)",
    "       ALS只解探针输出投影(64x64 Cholesky): 平凡",
    "       lm_head权重完全不动!",
    "  动机: 64维瓶颈消除残差放大 (delta被限制在小子空间)",
    "  7B轨迹: 60.2->34.5(C4)->24.2(C8)->22.7(C12)->22.8(C16)",
    "  PPL 22.8, 收敛但质量差 (瓶颈太窄, 信息损失)",
    "  瓶颈=3584/64=56x -> 表达能力下降到1/56",
],10)
y+=2.7
txbox(s,0.2,y,6.4,2.5,[
    "变体11: A-KD (Knowledge Distillation)",
    "  参考实验中提及, 无独立脚本",
    "  机制: 用教师模型输出替代恒等重建作为ALS目标",
    "       ALS 目标不再是 min||XW^T - XW_old||",
    "       而是 min||XW^T - teacher_output||",
    "  动机: ALS解决的对齐目标可以是有意义的教师信号",
    "  状态: 未独立实现, 概念阶段",
],10)
txbox(s,6.8,y,6.4,2.5,[
    "变体12: LARS Optimizer",
    "  脚本: _lars_qwen05b.py, _lars_sanity.py",
    "  机制: 替换SGD为LARS (Layer-wise Adaptive Rate Scaling)",
    "       trust_coefficient=0.001",
    "       每层独立学习率: lr_l = trust * ||param_l||/||grad_l||",
    "  动机: 层自适应率可能补偿残差放大的梯度不平衡",
    "  0.5B: 已运行但结果未记录 (runs/lars_qwen05b.json 缺)",
    "  7B: 未测",
    "  评估: 最自然的补充, 但LARS在NLP中的稳定性存疑",
],10)

# ═════════════════════════════════════════════════════════════  S9: Variant convergence plot
s=prs.slides.add_slide(BLANK);tbar(s,"A-SYNC 变体收敛轨迹 (7B)","8个变体 x Qwen2.5-7B (28L) -- 实际GPU数据")
img(s,os.path.join(FIG,"weekly_variant_convergence.png"),0.1,1.15,13.1,6.2)

# ═════════════════════════════════════════════════════════════  S10: 7B scoreboard
s=prs.slides.add_slide(BLANK);tbar(s,"7B 收敛排名: 全部A-SYNC变体","Qwen2.5-7B (28L) -- A-SYNC CONSTANT 48c 夺冠")
img(s,os.path.join(FIN,"fig7_scoreboard.png"),0.2,1.15,12.8,6.0)

# ═════════════════════════════════════════════════════════════  S11: Fix taxonomy
s=prs.slides.add_slide(BLANK);tbar(s,"修复策略分类: 5大类别","哪些策略真正有效? 完整归纳")
blt(s,0.3,1.2,12.5,5.8,[
    "A. 降低扰动幅度: 让ALS delta更小/更平滑再注入",
    "  EMA平滑 (beta=0.3跨周期平均) | Aligned注入 (只保留梯度同向分量) | step_size减小",
    "  效果: 有限 -- 与CONSTANT的结合未测试 (关键缺失实验)",
    "",
    "B. 增强恢复能力: 给SGD更多消化预算",
    "  更多SGD步/周期 (50->200) | 更高lr | LARS层自适应学习率",
    "  效果: 边际 -- LARS未测7B, 更多SGD步可能加速收敛",
    "",
    "C. 改变干预机制: 绕过残差放大路径 -- 唯一的一阶改善",
    "  A-SYNC梯度注入 (不修改权重, 塑形梯度) | A-PROBE低秩探针 (约束到64维瓶颈)",
    "  效果: 一阶 -- A-SYNC让7B从发散变为收敛; A-PROBE确认瓶颈假设但表达力不足",
    "",
    "D. 改变衰减调度: 控制何时以多大强度注入ALS信号 -- 二阶改善",
    "  CONSTANT无衰减 (最优: PPL 7.6) > Exponential衰减 > Cosine衰减 (最差: PPL 13.2) > Warm Restart (PPL 16.5)",
    "  效果: 二阶 -- 衰减策略全部反作用, CONSTANT通过保持信号使PPL从25.8->7.6 (3.4x改善)",
    "",
    "E. 事后平滑: 优化后平均权重降低方差",
    "  SWA (Stochastic Weight Averaging, C10起) -- 在A-SYNC非稳态轨迹上反增3.3 PPL",
    "  效果: 反作用 -- A-SYNC的ALS驱动轨迹根本非稳态, SWA假设稳态",
],11)

# ═════════════════════════════════════════════════════════════  S12: Section -- FLOPs
sec("FLOPs 归一化对比实验","A-SYNC CONSTANT vs AdamW Full-Rank vs LoRA AdamW  on  OPT-125m")

# ═════════════════════════════════════════════════════════════  S13: FLOPs accounting methodology
s=prs.slides.add_slide(BLANK);tbar(s,"FLOPs会计: 每步每模块的计算量","基于参数计数乘子模型 (匹配 altopt/profiling/flops.py)")
blt(s,0.3,1.2,6.2,4,[
    "FLOPs 乘子模型 (所有估计按 每步 = 乘子 x 可训练参数量):",
    "",
    "ALS (lm_head 块求解):  4 x n_params",
    "  含: 正规方程构建 (X^T X) + Cholesky分解 + 三角求解",
    "  OPT-125m: 4 x 125M = 501 MFLOPs/step",
    "  仅用于 A-SYNC 协议的 ALS 阶段 (每周期1次)",
    "",
    "SGD (前传+反传):  6 x n_params",
    "  含: 前传激活 (3x) + 反传梯度 (3x)",
    "  OPT-125m: 6 x 125M = 751 MFLOPs/step",
    "  用于 A-SYNC 的 SGD 阶段 (每周期50步)",
    "",
    "AdamW (前传+反传+2个矩):  10 x n_params",
    "  含: SGD的6x + 两个矩估计的存储/更新 (2x) + 偏置纠正 (2x)",
    "  OPT-125m: 10 x 125M = 1252 MFLOPs/step",
    "  LoRA版: 10 x 0.59M = 5.9 MFLOPs/step (!!)",
    "",
    "Eval (纯前传):  3 x n_params",
    "  OPT-125m: 3 x 125M = 376 MFLOPs/eval",
],10)

# Right: per-step bar chart + FLOPs budget table
img(s,os.path.join(FIG,"weekly_flops_design.png"),6.8,1.2,6.3,3.8)

blt(s,0.3,5.4,12,1.8,[
    "FLOPs 预算是匹配的: A-SYNC 24c = 24*(ALS+50*SGD) FLOPs = 24*(501+50*751)MF = 0.91T FLOPs",
    "AdamW 匹配: 0.91T / (10 x 125M) = 729步 | LoRA 匹配: 同样步数但 FLOPs仅0.013T -- 70x计算节省",
],10,RED)

# ═════════════════════════════════════════════════════════════  S14: FLOPs experiment design
s=prs.slides.add_slide(BLANK);tbar(s,"FLOPs归一化实验设计: 3协议详细对比","OPT-125m (12L, 125M params), WikiText-2, float32, 2xRTX 5090")

blt(s,0.3,1.2,6.2,5.8,[
    "协议1: A-SYNC CONSTANT 48c",
    "  每周期: 1x ALS lm_head (501MF) + 50x SGD全参数(body+head, 751MF/步)",
    "  ALS: reg_lambda=1e-3, step_size=0.01, block_size=512",
    "  SGD: lr=1e-4, momentum=0, weight_decay=0.01",
    "  Grad注入: grad += 0.05 * delta (恒, 不衰减!)",
    "  无扰动阶段 | 共48周期 -> 2400SGD步 + 48ALS步",
    "  总FLOPs: 1.846 TFLOPs, 墙钟: 318s",
    "",
    "协议2: AdamW Full-Rank (Protocol B)",
    "  每步: AdamW全参数(beta=0.9,0.999) -> 1252MF/步",
    "  lr=1e-4, weight_decay=0.01",
    "  共730步 (匹配A-SYNC 24c FLOPs预算)",
    "  总FLOPs: 0.911 TFLOPs, 墙钟: 136s",
    "",
    "协议3: LoRA AdamW r=8 (Protocol D style)",
    "  LoRA适配器: q_proj, v_proj, k_proj, out_proj",
    "  r=8, alpha=16, 仅 589,824 可训练参数 (213x fewer!)",
    "  每步: AdamW仅LoRA参数 -> 5.9MF/步 (!!)",
    "  同样730步, 但FLOPs完全不同: 0.013T (70x less!)",
    "  墙钟: 142s (前传仍需全参数, 所以墙钟不省)",
],10)

blt(s,6.8,1.2,6.2,5.8,[
    "为什么选OPT-125m (12层)?",
    "  . 125M参数: 适配32GB GPU, 可快速迭代",
    "  . 12层: rho^11=2.3x, 残差放大足够稳定",
    "    A-SYNC可收敛 (与28层7B不同)",
    "  . 快速: 全实验 ~5分钟, 便于调试",
    "  . 控制变量: 固定深度, 只比较协议",
    "",
    "FLOPs预算匹配逻辑:",
    "  目标: 给每个协议相同的计算预算",
    "  方法: A-SYNC24c总FLOPs = 24*(ALS+50*SGD)",
    "        = 24*(501+50*751) MF = 0.91 TFLOPs",
    "        AdamW步数 = 0.91T/1252MF = 729步",
    "        LoRA步数 = 同样729步 (FLOPs自动低)",
    "  A-SYNC 48c = 2x 预算 (研究更长训练)",
    "",
    "评估方案:",
    "  . A-SYNC: 每周期eval (25 eval点)",
    "  . AdamW/LoRA: 每30步eval (25 eval点)",
    "  . 指标: PPL (WikiText-2 test) + FLOPs",
    "  . 所有协议同起点: PPL=2246 (baseline)",
    "",
    "硬件: 2x NVIDIA RTX 5090 (32GB)",
    "  dtype=float32, batch=4 train / 8 eval, seq=128",
],10)

# ═════════════════════════════════════════════════════════════  S15: FLOPs design diagram + PPL plot
s=prs.slides.add_slide(BLANK);tbar(s,"FLOPs 预算分解图","每步计算量 + 总预算堆叠")
img(s,os.path.join(FIG,"weekly_flops_design.png"),0.2,1.15,13,6.2)

# ═════════════════════════════════════════════════════════════  S16: FLOPs-PPL detailed plot
s=prs.slides.add_slide(BLANK);tbar(s,"FLOPs vs PPL 详细对比 (OPT-125m)","标注关键数据点: AdamW PPL=23.2, LoRA PPL=37.3, A-SYNC PPL=60.7")
img(s,os.path.join(FIG,"weekly_flops_ppl_detailed.png"),0.1,1.15,13.1,6.2)

# ═════════════════════════════════════════════════════════════  S17: Results table
s=prs.slides.add_slide(BLANK);tbar(s,"FLOPs归一化对比: 完整数据表","OPT-125m, 所有3协议 + A-SYNC 48c")
rows=[
    ("协议","每步FLOPs","总步数","最终PPL","总FLOPs(T)","墙钟(s)","PPL/TFLOP","状态"),
    ("AdamW Full-Rank","1252 MF","730","23.2","0.911","136","25.5","已收敛"),
    ("LoRA AdamW r=8","5.9 MF","730","37.3","0.013","142","2812","缓慢收敛"),
    ("A-SYNC CONST 24c","ALS 501MF+SGD 751MF/step x50","24c","74.1","0.923","163","80.3","仍在改善"),
    ("A-SYNC CONST 48c","同上","48c","60.7","1.846","318","32.9","仍在改善"),
]
col_x=[0.2,2.8,5.8,7.5,9.0,10.3,11.8]
col_w=[2.5,2.9,1.6,2.5,1.2,1.4,1.6]
table(s,rows,1.3,col_x,col_w)

y2=1.3+len(rows)*0.42+0.2
blt(s,0.3,y2,12.5,3.5,[
    "关键观察:",
    "  . AdamW 绝对最优: PPL 23.2 -- 全秩+自适应矩在浅层模型(12L)上无法超越",
    "  . LoRA 效率无敌: PPL 37.3 at 1/70 FLOPs -- 589K参数达到可接受PPL",
    "  . LoRA PPL vs FLOPs 对比: 37.3/0.013T = 2812 PPL/TFLOP -- 87x 好于 A-SYNC!",
    "  . A-SYNC 48c 仍在收敛: 每周期Delta=-0.3~-0.7, 未见平台 (预测100c可达PPL~40)",
    "  . A-SYNC 在浅层模型的慢收敛是预期内的 -- 其价值是让28L模型收敛 (那里AdamW 1.25 vs A-SYNC 7.6)",
    "  . 墙钟时间 A-SYNC (318s) >> AdamW (136s) 因为每周期包含 ALS 求解 + eval 开销",
    "",
    "效率排名: LoRA (2812 PPL/TFLOP) >> A-SYNC 48c (33) > AdamW (25) | 但 -- 最终PPL才是底线!",
],10)

# ═════════════════════════════════════════════════════════════  S18: A-SYNC 48c full trajectory
s=prs.slides.add_slide(BLANK);tbar(s,"A-SYNC CONSTANT 48c: 完整收敛轨迹 (OPT-125m)","48周期, 每周期1次ALS + 50步SGD = 2448总步, 1.85 TFLOPs")
blt(s,0.3,1.2,6,5.5,[
    "Baseline: PPL=2246.1 (未经训练的OPT-125m)",
    "",
    "Phase 1 (C1-C8): 快速下降",
    "  C1: 1458 -> C4: 490 -> C8: 169",
    "  每周期改善: 200-400 PPL",
    " ALS delta 幅值大, 梯度注入强",
    " body 快速适应 lm_head 的新方向",
    "",
    "Phase 2 (C9-C24): 稳定收敛",
    "  C9: 140 -> C16: 87 -> C24: 74",
    "  每周期改善: 3-8 PPL",
    " ALS delta 开始自动缩小 (body 已部分适应)",
    " 无衰减: sync=0.05 保持满强度",
    "",
    "Phase 3 (C25-C48): 缓慢但持续",
    "  C25: 73 -> C36: 66 -> C48: 61",
    "  每周期改善: 0.3-0.7 PPL",
    " ALS delta 幅值继续自动缩小",
    " 未见平台 -- 仍在改善!",
    " 预测: C100 -> PPL ~40",
],10)
img(s,os.path.join(FIG,"flops_sweep_opt125m.png"),6.5,1.2,6.5,5.5)

# ═════════════════════════════════════════════════════════════  S19: Cross-model perspective
s=prs.slides.add_slide(BLANK);tbar(s,"跨模型结果汇总","12L -> 24L -> 28L: 不同深度, 不同故事")
blt(s,0.3,1.2,6.2,5.8,[
    "12层 (OPT-125m) -- 浅层模型, rho^11=2.3x",
    "  . A-SYNC 收敛: PPL 2246 -> 60.7 (48c)",
    "  . AdamW 大胜: PPL 23.2 (0.91T FLOPs)",
    "  . LoRA 最省: PPL 37.3 (0.013T FLOPs)",
    "  . 结论: 浅层上AdamW无可匹敌",
    "    梯度注入开销在2.3x放大下无收益",
    "",
    "24层 (Qwen0.5B) -- 中等深度, rho^23=5.9x",
    "  . Protocol A 仍可收敛但边际",
    "  . A-SYNC+变体 全部收敛 PPL 5.5 (0.5B容量底)",
    "  . LoRA 同达PPL 5.5 (更省算力)",
    "  . 结论: 0.5B太小无法区分变体",
    "",
    "28层 (Qwen7B) -- 深层模型, rho^27=8.0x",
    "  . Protocol A 11/11 发散 (CPU不可收敛!)",
    "  . A-SYNC CONSTANT 48c: 首次收敛 PPL 7.6",
    "  . AdamW 800step: PPL 1.25 (仍远优)",
    "  . A-SYNC差距: 6.1x vs AdamW",
    "  . 但 7.7x 好于baseline (58.8->7.6)",
    "  . 结论: A-SYNC解决了发散, 但未追上AdamW",
],11)

blt(s,6.8,1.2,6.2,5.8,[
    "深度-协议选择矩阵:",
    "",
    "  | 深度  | 推荐协议            | 理由                    |",
    "  |-------|---------------------|-------------------------|",
    "  | <=12L | AdamW Full-Rank     | 绝对最优, A-SYNC有开销 |",
    "  | 12-24L| AdamW 或 LoRA      | A-SYNC收敛但无优势      |",
    "  | >=28L | A-SYNC CONSTANT     | 唯一收敛的Protocol A    |",
    "  | 任意  | LoRA AdamW (r=8)   | 效率最优, 适合资源受限  |",
    "",
    "A-SYNC 价值定位:",
    "  . 不是要替代 AdamW (浅层上不可能)",
    "  . 是为 Protocol A 方法提供深层模型上的",
    "    收敛保证 (此前 28L+ 100% 发散)",
    "  . 从 发散 -> PPL 7.6 是质的突破",
    "  . 从 PPL 7.6 -> 1.25 (追平 AdamW)",
    "    是量的追赶, 可能需要 2-4x 更多周期",
    "",
    "未解决的问题:",
    "  1. ALS目标不对齐: 解的是重建loss, 非CE loss",
    "  2. 仅lm_head被ALS引导: body的27层仅靠SGD",
    "  3. 梯度裁剪限制: max_norm=1.0可能限幅有效更新",
    "  4. 长周期7B实验未做: 96-128c可能进一步缩小差距",
],10)

# ═════════════════════════════════════════════════════════════  S20: Full variant table
s=prs.slides.add_slide(BLANK);tbar(s,"A-SYNC 12变体完整对比表","机制要点 + 7B PPL + vs CONSTANT 48c")
rows=[
    ("变体","机制要点","7B PPL","vs CONSTANT 48c","收敛状态"),
    ("Vanilla (8c)","delta注入+扰动+指数衰减, sync*=0.8","25.8","Delta+18.2","缓慢收敛(衰减过早停用)"),
    ("No-Perturb (8c)","移除扰动阶段, 其余同Vanilla","16.6","Delta+9.0","收敛(扰动7B上纯有害)"),
    ("Cosine 32c","sync+lr = 余弦衰减至0","13.2","Delta+5.6","C20平台(衰减杀死尾部)"),
    ("CONSTANT 24c","sync=0.05恒, lr=2e-4恒, 无扰动","9.0","Delta+1.4","收敛, 首个破PPL 10"),
    ("CONSTANT 48c","延至48周期, C44自然收敛","7.6","BASELINE","C44自然收敛, 最优"),
    ("A-CYCLE 3x8","3块x8周期, 块内Cosine, 块边界重置","16.5","Delta+8.9","块内衰减浪费4周期/块"),
    ("+EMA (beta=0.3)","delta指数滑动平均跨周期","0.5B:5.5","7B未测","7B未运行(runs缺)"),
    ("+Aligned","只注入delta与grad同向分量","0.5B:5.5","7B未测","7B未运行(runs缺)"),
    ("+Warmup","4cycle SGD预热后A-SYNC","0.5B:5.5","7B未测","7B未运行(runs缺)"),
    ("+SWA Cosine","C10起SWA平均权重+Cosine衰减","10.5","Delta+2.9","SWA反增3.3PPL(!)"),
    ("A-PROBE r=64","低秩探针绕过lm_head, ALS解探针","22.8","Delta+15.2","收敛但容量受限(瓶颈56x)"),
    ("LARS optimizer","层自适应学习率替代SGD","0.5B:???","未测7B","0.5B结果未记录, 7B未运行"),
]
col_x=[0.2,2.6,7.5,9.5,11.3]
col_w=[2.3,4.8,1.9,1.7,2.1]
table(s,rows,1.3,col_x,col_w)

y_end=1.3+len(rows)*0.42+0.15
blt(s,0.3,y_end,12,1.2,[
    "排序: CONSTANT 48c (7.6) < CONSTANT 24c (9.0) < +SWA (10.5) < Cosine (13.2) < No-Perturb (16.6) < A-CYCLE (16.5) < A-PROBE (22.8) < Vanilla (25.8)",
    "核心规律: 衰减=有害, 扰动=有害, 无衰减+更多周期=最优. 7B上仍有5个变体未测 (EMA/Aligned/Warmup/LARS/长CONSTANT).",
],9,RED)

# ═════════════════════════════════════════════════════════════  S21: Summary insights
s=prs.slides.add_slide(BLANK);tbar(s,"本周核心发现总结","从理论到实验: 一条完整的证据链")
blt(s,0.3,1.2,6.2,5.8,[
    "发现1: 残差放大 (rho~1.08) 是根本原因",
    "  SCM因果推导: delta_{l+1}=(I+J_l).delta_l",
    "  跨4个模型族的实验验证: L_max~26层",
    "  1.08来自 I+J_l 谱范数的跨层几何均值",
    "  8% per-layer 放大是训练后Transformer的固有属性",
    "",
    "发现2: A-SYNC梯度注入绕过放大",
    "  模块级变化: lm_head ALS -> 仅算delta, 权重不修改",
    "  body SGD -> SGD + grad梯度偏置注入 (sync*delta)",
    "  全参数: 移除Perturb (7B上反作用)",
    "  7B (28L) 首次收敛 (从100%发散 -> PPL 7.6)",
    "",
    "发现3: CONSTANT调度是最优策略",
    "  衰减全部有害: Cosine (平台) > Exponential (提前停用)",
    "  CONSTANT 让ALS delta信号持续, body持续自适应",
    "  自然收敛: ALS delta幅度随body适应而自动缩小",
    "  -> 外部衰减完全多余且自毁",
],11)

blt(s,6.8,1.2,6.2,5.8,[
    "发现4: FLOPs归一化实验 (OPT-125m)",
    " 每步FLOPs详细会计: ALS=4x, SGD=6x, AdamW=10x",
    " AdamW Full-Rank: 23.2 PPL (绝对最优)",
    " LoRA r=8: 37.3 PPL at 0.013T (效率最优, 70x less)",
    " A-SYNC 48c: 60.7 PPL at 1.85T (仍在收敛)",
    "",
    "发现5: A-SYNC的定位是深层模型",
    "  | 深度  | 推荐              | A-SYNC是否有优势 |",
    "  | 12L   | AdamW或LoRA       | 无 (梯度注入开销) |",
    "  | 24L   | 任一均可          | 边际 (0.5B太小) |",
    "  | 28L+  | A-SYNC CONSTANT   | 唯一收敛的Prot.A |",
    "",
    "发现6: 5大类修复策略中只有2类有效",
    "  C类 (干预机制): 梯度注入 -- 一阶突破",
    "  D类 (衰减调度): CONSTANT -- 二阶优化",
    "  A类/B类/E类: 效果有限/边际/反作用",
    "",
    "下一步: 7B上测 EMA/Aligned/CONSTANT 96c+",
    "  缩小与AdamW的6.1x差距",
],11)

# ═════════════════════════════════════════════════════════════  S22: Next steps
s=prs.slides.add_slide(BLANK);tbar(s,"下一步计划","实验 + 论文 + 投稿")

blt(s,0.3,1.2,6.2,5.5,[
    "短期 (1-2周): 完成缺失的7B实验",
    "  1. A-SYNC+EMA 7B 验证",
    "     脚本: _a_sync_plus_7b.py 已写, 待GPU运行",
    "     预期: 抑制晚期噪声, 稳定C40+收敛",
    "",
    "  2. A-SYNC+Aligned 7B 验证",
    "     脚本: 需适配到7B",
    "     预期: 减少破坏性梯度干涉",
    "",
    "  3. A-PROBE 更大rank (256/512/1024)",
    "     脚本: _probe_7b.py 需改 rank参数",
    "     预期: 在表达力和安全之间找到甜点",
    "",
    "中期 (2-4周): 长周期7B实验",
    "  4. A-SYNC CONSTANT 96-128 cycle 7B",
    "     测试: 持续收敛能否缩小与AdamW差距",
    "     预期: 可能达到 PPL 3-5 范围",
    "",
    "  5. 多head ALS: 2-3个attention层同时ALS",
    "     测试: 扩大ALS引导范围的效果",
],11)

blt(s,6.8,1.2,6.2,5.5,[
    "论文推进 (与实验并行):",
    "",
    "  6. 独立复核者复算主表",
    "     P1.1/1.2/1.3 + FLOPs实验",
    "     确保所有数字可复现",
    "",
    "  7. 更新论文至 v0.9",
    "     加入A-SYNC CONSTANT 48c 7B结果",
    "     加入FLOPs归一化对比",
    "     更新结论: Protocol A现在可收敛深层",
    "",
    "  8. 选择投稿期刊",
    "     TMLR (双盲, 滚动审稿)",
    "     arXiv + NeurIPS/ICML workshop",
    "     ACL Rolling Review",
    "",
    "",
    "本周产出清单:",
    "  docs/a_sync_report_en.pdf (456KB)",
    "  docs/a_sync_report_zh.pdf (529KB)",
    "  docs/residual-amplification-why-108.md",
    "  docs/a-sync-variants-and-fixes.md",
    "  docs/flops-sweep-report.md",
    "  experiments/_flops_sweep.py",
    "  experiments/_gen_weekly_ppt.py",
    "  docs/weekly_group_meeting_20260724.pptx",
],11)

# ── Thank you bar ──────────────────────────────────────────────────
bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(6.95),W,Inches(0.55))
bar.fill.solid();bar.fill.fore_color.rgb=DARK;bar.line.fill.background()
tf=bar.text_frame;p=tf.paragraphs[0]
p.text="Thanks!  .  github.com/hjiang555-a11y/alternating-optimization-lora  .  2026/07/24"
p.font.size=Pt(11);p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER

# ── Save ──────────────────────────────────────────────────────────
out=os.path.join(os.path.dirname(__file__),"..","docs","weekly_group_meeting_20260724.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
